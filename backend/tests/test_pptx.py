import os

from pptx import Presentation
from pptx.util import Inches

from app.ingestion.pptx import PptxParser


def _make_fixture(path):
    prs = Presentation()
    for text in ["Intro to Algebra", "Solving for x"]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text_frame.text = text
    prs.save(path)


def test_pptx_extracts_text_per_slide(tmp_path):
    path = os.path.join(tmp_path, "deck.pptx")
    _make_fixture(path)
    parsed = PptxParser().parse(path)
    assert len(parsed.segments) == 2
    assert "Algebra" in parsed.segments[0][0]
    assert parsed.segments[0][1] == {"type": "pptx", "slide": 1}
    assert parsed.segments[1][1]["slide"] == 2
