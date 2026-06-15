import os

from pptx import Presentation

from app.ingestion.base import ParsedSource


class PptxParser:
    def parse(self, ref: str) -> ParsedSource:
        prs = Presentation(ref)
        segments: list[tuple[str, dict]] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
            joined = "\n".join(t for t in texts if t.strip())
            if joined.strip():
                segments.append((joined, {"type": "pptx", "slide": i}))
        if not segments:
            raise ValueError("PPTX has no extractable text")
        return ParsedSource(title=os.path.basename(ref), segments=segments)
